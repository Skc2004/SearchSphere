import PyPDF2
import os
from docx import Document
from pptx import Presentation
import pandas as pd
import pathlib
from datetime import datetime
import markdown
from bs4 import BeautifulSoup
from tqdm import tqdm

def get_meta(file_path: os.path) -> dict:
    """
    Function to get meta data of file
    args:
        file_path: os.path
    return
        dict
    """
    file_name = os.path.basename(file_path)
    file_type = pathlib.Path(file_path).suffix
    folder_location = os.path.dirname(file_path)

    if os.name == 'nt':#for windows
        create_time = os.path.getctime(file_path)

    else:
        stat = os.stat(file_path)
        create_time = stat.st_mtime

    create_date = datetime.fromtimestamp(create_time).strftime('%Y-%m-%d %H:%M:%S')

    return {
        "file_name" : file_name,
        "file_type" : file_type,
        "file_path" : file_path,
        "creation_date" : create_date
    }

def pdf_extractor(file_path: os.path) -> str:
    """
    extracts text from pdf
    Args:
        file_path : pdf file path
    Returns:
        text
    """ 
    with open(file_path , 'rb') as file:
        #pdf reader pointer
        reader = PyPDF2.PdfReader(file)
        # all pdf texts
        pdf_text = ''
        
        #if book: look for toc as it contains the most amount of information
        start_page = 0
        toc_start_page = None
        end_page = min(20 , len(reader.pages))

        #look for toc
        toc_keywords = ["table of content" , "contents" , "content" , "index" , "navigating"]
        
        #look for table of content
        for page_num , page in enumerate(reader.pages):

            text = page.extract_text().lower()

            if any(keyword in text for keyword in toc_keywords):
                toc_start_page = page_num

        #book traversal
        if toc_start_page is not None:
            first_page = reader.pages[0] if reader.pages[0] != '' else reader.pages[1]
            pdf_text += first_page.extract_text()

            for page_num in tqdm(range(toc_start_page , end_page)):
                page = reader.pages[page_num]
                pdf_text += page.extract_text()

                #break point -> 10000 words
                if len(pdf_text) >= 10000:
                    break
    
        else: #if toc not available
            for page_num in tqdm(range(start_page , end_page)):
                page = reader.pages[page_num]
                pdf_text += page.extract_text()
                if len(pdf_text) >= 10000:
                    break

    return pdf_text
            


def text_extractor(file_path: os.path)->str:
    """
    Function to extract text data from txt files
    args:
        file_path: os.path
    return
        text: str
    """
    with open(file_path , 'r') as file:
        text = file.read()
        #10000 words limit
        text = text[:10000]
        return text


def docs_extractor(file_path: os.path) -> str:
    doc = Document(file_path)
    txt = ''
    for para in tqdm(doc.paragraphs):
        txt += para.text + '\n'
        if len(txt) >= 10000:
            break

    return txt

def ppt_extractor(file_path: os.path) -> str:
    prs = Presentation(file_path)

    text = ""
    for slide in tqdm(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape , "text"):
                text += shape.text

    return text                


def excel_extractor(file_path: os.path) -> str:
    excel_file = pd.read_excel(file_path)
    
    text = ""
    col_name = excel_file.columns.to_list()
    if col_name is not None:
        for cols in tqdm(col_name):   
            text += cols

    else:
        #rely on meta data
        text = None
  

def markdown_extractor(file_path: os.path) -> str:
    with open(file_path ,  'r' , encoding='utf-8') as file:
        md_content = file.read()

    html_cont = markdown.markdown(md_content)

    soup = BeautifulSoup(html_cont , 'html.parser')
    text = soup.get_text()

    return text[:10000 if len(text) > 10000 else len(text)]




     








        
